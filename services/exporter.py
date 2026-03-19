import asyncio
import io
import logging
import math
import re
import zipfile

import httpx
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN

from services.downloader import fetch_bytes

logger = logging.getLogger(__name__)

HEADERS = {
    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 "
                  "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
}


def _sanitize_filename(name: str) -> str:
    name = re.sub(r"[^\w\s\-]", "", name)
    name = re.sub(r"\s+", "_", name.strip())
    return name[:60] or "logo"


def _detect_extension(content_type: str, image_bytes: bytes) -> str:
    ct = (content_type or "").lower()
    if "png" in ct:
        return ".png"
    if "jpeg" in ct or "jpg" in ct:
        return ".jpg"
    if "webp" in ct:
        return ".webp"
    if "svg" in ct:
        return ".svg"
    if "ico" in ct:
        return ".ico"
    # Fallback: check magic bytes
    if image_bytes[:4] == b"\x89PNG":
        return ".png"
    if image_bytes[:2] in (b"\xff\xd8",):
        return ".jpg"
    if image_bytes[:4] == b"<svg" or b"<svg" in image_bytes[:100]:
        return ".svg"
    return ".png"


async def build_zip(selections: list[dict]) -> io.BytesIO:
    async with httpx.AsyncClient(headers=HEADERS) as client:
        tasks = [client.get(item["url"], timeout=10.0, follow_redirects=True) for item in selections]
        responses = await asyncio.gather(*tasks, return_exceptions=True)

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for item, resp in zip(selections, responses):
            if isinstance(resp, Exception) or resp.status_code != 200:
                continue
            ext = _detect_extension(resp.headers.get("content-type", ""), resp.content)
            filename = _sanitize_filename(item["company"]) + ext
            zf.writestr(filename, resp.content)

    buffer.seek(0)
    return buffer


def _is_svg(raw: bytes) -> bool:
    try:
        snippet = raw[:512].decode("utf-8", errors="ignore").lower()
        return "<svg" in snippet
    except Exception:
        return False


def _to_png_bytes(raw: bytes) -> tuple[bytes | None, float]:
    """
    Convert any image bytes to PNG (required by python-pptx).
    Returns (png_bytes, aspect_ratio_h_over_w). Returns (None, 1.0) on failure.
    Handles SVG via cairosvg, all other formats via Pillow.
    """
    # SVG → PNG via cairosvg
    if _is_svg(raw):
        try:
            import cairosvg
            png = cairosvg.svg2png(bytestring=raw, output_width=500)
            with Image.open(io.BytesIO(png)) as img:
                ratio = img.height / img.width if img.width > 0 else 1.0
            return png, ratio
        except Exception:
            return None, 1.0

    # Raster images → PNG via Pillow
    try:
        with Image.open(io.BytesIO(raw)) as img:
            iw, ih = img.size
            ratio = ih / iw if iw > 0 else 1.0
            if img.mode not in ("RGB", "RGBA"):
                img = img.convert("RGBA")
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue(), ratio
    except Exception:
        return None, 1.0


async def _fetch_one(client: httpx.AsyncClient, item: dict) -> tuple[dict, bytes | None]:
    """Download a single logo URL, return (item, raw_bytes_or_None)."""
    url = item["url"]
    if url.startswith("data:"):
        from services.downloader import _decode_data_uri
        return item, _decode_data_uri(url)
    try:
        resp = await client.get(url, timeout=10.0, follow_redirects=True)
        if resp.status_code == 200:
            return item, resp.content
    except Exception as e:
        logger.warning("PPTX: download error for %s: %s", item["company"], e)
    return item, None


async def build_pptx(selections: list[dict], logo_width_cm: float = 5.0,
                     show_labels: bool = True) -> io.BytesIO:
    # Download all images in parallel
    async with httpx.AsyncClient(headers=HEADERS) as client:
        tasks = [_fetch_one(client, item) for item in selections]
        results = await asyncio.gather(*tasks)

    items_with_png: list[tuple[dict, bytes, float]] = []
    for item, raw in results:
        if not raw:
            logger.warning("PPTX: fetch failed for %s", item["company"])
            continue
        logger.info("PPTX: got %d bytes for %s, converting to PNG", len(raw), item["company"])
        png_bytes, ratio = _to_png_bytes(raw)
        if png_bytes:
            logger.info("PPTX: converted %s → %d PNG bytes (ratio=%.2f)", item["company"], len(png_bytes), ratio)
            items_with_png.append((item, png_bytes, ratio))
        else:
            logger.warning("PPTX: PNG conversion failed for %s", item["company"])

    logger.info("PPTX: building slide with %d / %d logos", len(items_with_png), len(selections))

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    n = len(items_with_png)
    if n == 0:
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    # Grid: pick columns so cells are roughly square
    cols = math.ceil(math.sqrt(n))
    rows = math.ceil(n / cols)

    # Generous outer margins — leave room for slide "border/frame" space
    margin_x = Cm(2.5)
    margin_y = Cm(2.0)
    label_h = Cm(0.75) if show_labels else 0
    cell_gap = Cm(0.4)   # gap between cells

    usable_w = prs.slide_width - margin_x * 2
    usable_h = prs.slide_height - margin_y * 2
    cell_w = int((usable_w - cell_gap * (cols - 1)) / cols)
    cell_h = int((usable_h - cell_gap * (rows - 1)) / rows)

    for idx, (item, png_bytes, ratio) in enumerate(items_with_png):
        col = idx % cols
        row = idx // cols

        cell_left = int(margin_x + col * (cell_w + cell_gap))
        cell_top = int(margin_y + row * (cell_h + cell_gap))

        # Company name label at bottom of cell (optional)
        if show_labels:
            try:
                lbl = slide.shapes.add_textbox(
                    cell_left, cell_top + cell_h - int(label_h), cell_w, int(label_h)
                )
                tf = lbl.text_frame
                tf.word_wrap = True
                tf.text = item["company"]
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                para.runs[0].font.size = Pt(9)
            except Exception as e:
                logger.warning("PPTX: label failed for %s: %s", item["company"], e)

        # Logo image — centered within cell, above label
        try:
            img_area_h = cell_h - int(label_h) - Cm(0.15)
            img_area_w = cell_w - Cm(0.3)

            logo_w = min(Cm(logo_width_cm), int(img_area_w))
            logo_h = int(logo_w * ratio)
            if logo_h > img_area_h:
                logo_h = int(img_area_h)
                logo_w = int(logo_h / ratio) if ratio > 0 else logo_w

            logo_w = max(logo_w, Cm(0.5))
            logo_h = max(logo_h, Cm(0.5))

            logo_left = int(cell_left + (cell_w - logo_w) / 2)
            logo_top = int(cell_top + (img_area_h - logo_h) / 2)

            slide.shapes.add_picture(
                io.BytesIO(png_bytes), logo_left, logo_top, logo_w, logo_h
            )
            logger.info("PPTX: added %s at col=%d row=%d", item["company"], col, row)
        except Exception as e:
            logger.warning("PPTX: add_picture failed for %s: %s", item["company"], e)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    logger.info("PPTX: done, size=%d bytes", len(output.getvalue()))
    return output
